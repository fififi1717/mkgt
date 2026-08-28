#!/usr/bin/env python3
"""
Affilior — Contrôle des graphiques natifs (donut, slide 6 "Allocation").

Constat crash-test réel (28/08/2026) : `assemble.py` et `finalize_deck.py`
ne touchent jamais aux graphiques natifs PowerPoint (chart1.xml / chart2.xml
du Canvas Master) — seuls les placeholders texte et tableaux sont remplis.
Résultat : un deck peut être livré avec les DEUX donuts encore aux valeurs
de démonstration du template (40/30/20/10 et 35/65) alors que tous les
textes autour affichent les vrais chiffres du client. `validate.py` ne le
détecte pas (XML structurellement valide) et `check_montants.py` ne le
détecte pas non plus (les valeurs de graphique ne sont pas des montants en
euros dans le texte des slides, ce sont des pourcentages dans un flux XML
de graphique séparé).

Comportement : NON BLOQUANT par défaut, même logique que check_montants.py
— on signale, on ne bloque jamais tout seul. La correction automatique
(--fix) n'est appliquée que si le dossier client fournit explicitement les
deux répartitions ET que les catégories correspondent exactement à celles
du template (aucune invention de catégorie, R1).

Usage :
    python3 check_charts.py deck.pptx --dossier dossier_client.json
    python3 check_charts.py deck.pptx --dossier dossier_client.json --fix --out deck_fixed.pptx

dossier_client.json doit contenir, en plus des clés existantes :
{
  "allocation_pct": {"Immobilier": 45, "Financier": 38, "Retraite": 12, "Liquidités": 5},
  "liquidite_pct": {"Liquide": 20, "Illiquide": 80}
}
"""
import argparse
import copy
import json
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData

# Signature des valeurs de démonstration livrées avec le Canvas Master.
# Si un graphique du deck final égale encore CETTE signature, c'est la
# preuve qu'il n'a jamais été mis à jour pour ce dossier — c'est le bug,
# pas une coïncidence plausible avec un vrai client.
DEMO_SIGNATURES = [
    (("Immobilier", "Financier", "Retraite", "Liquidités"), (40.0, 30.0, 20.0, 10.0)),
    (("Liquide", "Illiquide"), (35.0, 65.0)),
]


def _matches_demo(categories, values):
    cats = tuple(categories)
    vals = tuple(round(v, 1) for v in values)
    for demo_cats, demo_vals in DEMO_SIGNATURES:
        if cats == demo_cats and vals == demo_vals:
            return True
    return False


def check(pptx_path, dossier):
    """Retourne une liste de rapports : (slide_n, series_name, categories, values, statut)."""
    prs = Presentation(pptx_path)
    alerts = []
    fixes = []  # (chart_obj, CategoryChartData) à appliquer si --fix

    expected = {
        "Allocation": dossier.get("allocation_pct"),
        "Liquidité": dossier.get("liquidite_pct"),
    }

    for slide_n, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart = shape.chart
            plot = chart.plots[0]
            categories = list(plot.categories)
            for series in plot.series:
                values = list(series.values)
                is_demo = _matches_demo(categories, values)
                exp = expected.get(series.name)

                if is_demo and not exp:
                    alerts.append((slide_n, series.name, categories, values,
                                    "⚠ valeurs de démo du template, aucune donnée client fournie pour ce graphique"))
                    continue

                if is_demo and exp:
                    if set(exp.keys()) != set(categories):
                        alerts.append((slide_n, series.name, categories, values,
                                        f"⚠ valeurs de démo détectées mais catégories fournies ({list(exp.keys())}) "
                                        f"ne correspondent pas aux catégories du template ({categories}) — "
                                        f"pas de correction automatique possible sans risque d'erreur de mapping"))
                        continue
                    total = sum(exp.values())
                    if abs(total - 100) > 0.5:
                        alerts.append((slide_n, series.name, categories, values,
                                        f"⚠ valeurs de démo détectées, répartition fournie somme à {total}% (≠100%) — à corriger avant fix"))
                        continue
                    alerts.append((slide_n, series.name, categories, values,
                                    "⚠ valeurs de démo détectées — correction disponible (--fix)"))
                    new_values = [exp[c] for c in categories]
                    cd = CategoryChartData()
                    cd.categories = categories
                    cd.add_series(series.name, new_values)
                    fixes.append((chart, cd))
                    continue

                if not is_demo and exp:
                    computed = [round(exp.get(c, 0), 1) for c in categories]
                    if tuple(round(v, 1) for v in values) != tuple(computed):
                        alerts.append((slide_n, series.name, categories, values,
                                        f"⚠ graphique déjà modifié mais ne correspond pas à allocation_pct fourni "
                                        f"(attendu {computed}) — vérifier manuellement, pas de fix automatique "
                                        f"sur un graphique déjà personnalisé"))
                # not is_demo and not exp -> rien à signaler (déjà rempli, pas de donnée de contrôle fournie)

    return alerts, fixes


def apply_fixes(pptx_path, dossier, out_path):
    """CORRECTIF (déduplication, constat 29/08/2026) : implémentation UNIQUE de
    la correction des graphiques, partagée entre check_charts.py --fix et
    finalize_deck.py (qui dupliquait auparavant exactement cette logique en
    ligne, avec le risque de divergence silencieuse entre les deux versions).

    Ouvre pptx_path, corrige les séries non ambiguës (mêmes règles que check()),
    sauvegarde vers out_path. Retourne (applied, per_series) où per_series est
    une liste de tuples (slide_n, series_name, corrige: bool, raison) — assez
    détaillée pour qu'un message de synthèse distingue "corrigé" de "laissé
    tel quel" SÉRIE PAR SÉRIE, même quand plusieurs graphiques cohabitent sur
    la même slide (cf. constat 29/08/2026 : le message précédent parlait de
    "slide corrigée" sans préciser qu'un des deux donuts pouvait rester faux).
    """
    prs = Presentation(pptx_path)
    expected = {"Allocation": dossier.get("allocation_pct"),
                "Liquidité": dossier.get("liquidite_pct")}
    per_series = []
    applied = 0
    for slide_n, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart = shape.chart
            plot = chart.plots[0]
            categories = list(plot.categories)
            for series in plot.series:
                values = list(series.values)
                if not _matches_demo(categories, values):
                    continue  # déjà personnalisé ou non concerné, rien à corriger ici
                exp = expected.get(series.name)
                if not exp:
                    per_series.append((slide_n, series.name, False, "aucune donnée client fournie"))
                    continue
                if set(exp.keys()) != set(categories):
                    per_series.append((slide_n, series.name, False, "catégories fournies ≠ catégories du template"))
                    continue
                total = sum(exp.values())
                if abs(total - 100) > 0.5:
                    per_series.append((slide_n, series.name, False, f"somme fournie à {total}% (≠100%)"))
                    continue
                cd = CategoryChartData()
                cd.categories = categories
                cd.add_series(series.name, [exp[c] for c in categories])
                chart.replace_data(cd)
                applied += 1
                per_series.append((slide_n, series.name, True, "corrigé"))
    prs.save(out_path)
    return applied, per_series


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("src")
    ap.add_argument("--dossier", required=True)
    ap.add_argument("--fix", action="store_true",
                     help="Applique la correction quand elle est possible sans ambiguïté "
                          "(catégories identiques au template, somme à 100%%).")
    ap.add_argument("--out", help="Requis si --fix.")
    args = ap.parse_args()

    if args.fix and not args.out:
        print("Erreur : --fix nécessite --out.", file=sys.stderr)
        sys.exit(2)

    dossier = json.load(open(args.dossier, encoding="utf-8"))
    alerts, fixes = check(args.src, dossier)

    print(f"=== Contrôle graphiques natifs (donuts) — {args.src} ===")
    if not alerts:
        print("✓ Aucun graphique laissé aux valeurs de démonstration du template.")
    for slide_n, name, cats, vals, msg in alerts:
        print(f"  slide {slide_n:<3} série « {name} » {dict(zip(cats, vals))}")
        print(f"    {msg}")

    if args.fix:
        if not fixes:
            print("\nAucune correction applicable automatiquement (voir alertes ci-dessus).")
            return
        # CORRECTIF (déduplication, constat 29/08/2026) : appelle désormais
        # apply_fixes(), implémentation unique partagée avec finalize_deck.py.
        applied, per_series = apply_fixes(args.src, dossier, args.out)
        print(f"\n✓ {applied} graphique(s) corrigé(s) -> {args.out}")
        for slide_n, name, ok, reason in per_series:
            if not ok:
                print(f"  ⚠ slide {slide_n} « {name} » NON corrigé — {reason}")


if __name__ == "__main__":
    main()
