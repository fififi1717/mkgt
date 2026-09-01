#!/usr/bin/env python3
"""
personalize_canvas.py — Affilior, étape "personnalisation Canvas" (nouvelle,
avant assemble.py). Remplace le remplissage XML manuel fait par Claude à
chaque session Gate 5 — introduit après le crash-test du 30/08/2026 qui a mis
en évidence plusieurs bugs récurrents de ce remplissage manuel :
  - couleur bronze (8E5B3F) écrasée par erreur en gris corps au remplissage
  - labels % de la slide 6 (allocation/liquidité) jamais synchronisés avec
    le graphique natif que corrige check_charts.py
  - slide 5 (Patrimoine) : ancien tableau générique 10 lignes, jamais utilisé
    en pratique (la maquette réelle 2 colonnes/famille était documentée mais
    jamais scriptée)
  - slide 7 : pas de représentation visuelle du barème IR sous le TMI

Usage :
    python3 personalize_canvas.py --canvas Cans_Mstr.pptx --plan plan.json --out Cans_Mstr_perso.pptx
    python3 personalize_canvas.py --canvas Cans_Mstr.pptx --show-fields 7   # aperçu de l'ordre des champs, sans remplir

plan.json attend (toutes les clés optionnelles sauf 'replacements') :
{
  "mois_annee": "AOÛT 2026",
  "replacements": {"1": [...], "4": [...], "7": [...], "8": [...], "9": [...],
                    "10": [...], "12": [...], "13": [...], "14": [...],
                    "15": [...], "16": [...], "17": [...], "18": [...]},
  "allocation_pct": {"Immobilier": 64.0, "Financier": 26.0, "Retraite": 6.4, "Liquidités": 3.6},
  "liquidite_pct": {"Liquide": 36.0, "Illiquide": 64.0},
  "patrimoine": {
    "familles": [{"nom": "ASSURANCE VIE", "actifs": [{"nom": "...", "proprietaire": "...", "montant": 450000}]}],
    "total_net": 1250000
  },
  "tmi_label": "45 %"
}

Les valeurs de 'replacements' sont produites par content_plans.build_plan()
à partir des données structurées du dossier client (cf. content_plans.py).

CORRECTIF v4.20 (30/08/2026, crash-test double client) : la clé "17" (slide
Accompagnement, bloc "Engagements spécifiques à ce dossier") était absente
de ce docstring et de template_positions.json["slides_personnalisees"] —
elle sortait donc TOUJOURS avec les 6 placeholders visibles, sur tous les
dossiers générés jusqu'ici, sans qu'aucun script ne le signale comme un
défaut (le format "[Engagement 1]" n'est pas reconnu comme incomplet par
check_template_residue.py, cf. correctif séparé dans ce script). Attend
exactement 6 valeurs, dans cet ordre (vérifié par inspection XML directe
du Cans_Mstr.pptx en production) :
  [0] Intitulé engagement 1   [1] Détail engagement 1 (comment/fréquence/interlocuteur)
  [2] Intitulé engagement 2   [3] Détail engagement 2
  [4] Intitulé engagement 3   [5] Détail engagement 3
Cette section est éditoriale (ton, fréquence de suivi) — pas de montant
attendu, donc sans impact sur check_montants.py.
La clé "5" (Votre patrimoine) ne doit JAMAIS être fournie dans 'replacements' :
cette slide est reconstruite nativement via 'patrimoine', pas remplie par
substitution de texte (cf. rebuild_slide5_patrimoine).

CORRECTIF 01/09/2026 (reconstruction Canvas — retour consultant "cases jamais
voulues" + "ligne horizontale jamais implémentée") : le Cans_Mstr.pptx binaire
avait pris du retard sur deux décisions déjà actées dans le skill. Corrigé
en binaire ce jour ; les clés "4" et "12" changent de format :

  "4" (Situation) — ATTEND DÉSORMAIS EXACTEMENT 6 VALEURS, une phrase complète
  reformulée par puce (ex. "Vous avez 58 et 55 ans, mariés sous le régime de
  la communauté réduite aux acquêts.") — jamais un format libellé/valeur.
  Remplace l'ancien format à 8 valeurs (champs séparés). Si le dossier a moins
  de 6 items réels, remplir les puces excédentaires avec un texte de
  comblement (même logique que slides 12/13/15 : filler puis retrait par
  remove_unused_slots.py via nb_elements_reels["4"], jamais de puce vide
  laissée dans le livrable). Accords de genre obligatoires (Monsieur/Madame).

  "12" (Projets) — toujours 6 valeurs, mais le sens des 2 valeurs par ligne
  change : [0]=mot-clé court (1-3 mots, ex. "Transmission") [1]=phrase à la
  1ère personne. Remplace l'ancien sens [intitulé, description longue].
  Structurellement identique à avant (aucun changement de code nécessaire
  ici) — seul template_positions.json changent (format "groupé" par ligne :
  mot-clé + phrase + filet, cf. remove_unused_slots.py).
"""
import argparse
import json
import os
import re
import shutil
import zipfile

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Constantes de charte — cf. SKILL.md §5.4, corrigées le 30/08/2026 (2 passes) :
# 1re passe : R16 mentionnait C4B5A5 comme seule couleur de placeholder, alors
# que le Canvas avait déjà migré vers le bronze clair 8E5B3F — corrigé pour ne
# plus écraser 8E5B3F par du gris.
# 2e passe (ce jour, confirmée par Marc) : le bronze clair 8E5B3F est lui-même
# remplacé par le bronze fort 7A3E1D partout où il apparaît. Le beige C4B5A5
# n'est PAS concerné par cette migration — il reste une couleur distincte de
# la hiérarchie (coexiste avec le bronze fort), donc laissé tel quel, jamais
# converti automatiquement ici.
# ---------------------------------------------------------------------------
COLOR_CORPS = "4A5568"
COLOR_BRONZE_CLAIR = "8E5B3F"   # couleur de saisie du Canvas, migre vers une des 2 couleurs finales ci-dessous
COLOR_BRONZE_FORT = "7A3E1D"    # couleur finale par défaut (valeurs numériques / chiffres clés)
COLOR_TAUPE = "524A44"          # couleur finale pour les slides à contenu narratif/descriptif
PLACEHOLDER_COLORS_TO_MIGRATE = {COLOR_BRONZE_CLAIR, COLOR_BRONZE_FORT}
# Slides narratives (texte descriptif, pas de chiffre clé) -> taupe. Toutes les autres -> bronze fort.
TAUPE_SLIDES = {"4", "9", "12"}

def color_for_slide(slide_num):
    return COLOR_TAUPE if str(slide_num) in TAUPE_SLIDES else COLOR_BRONZE_FORT

NAVY = RGBColor(0x1C, 0x2B, 0x3A)
BORDEAUX = RGBColor(0x8B, 0x1A, 0x1A)
CREME = RGBColor(0xF3, 0xF0, 0xEB)
CORPS = RGBColor(0x4A, 0x55, 0x68)
BRONZE_FORT_RGB = RGBColor(0x7A, 0x3E, 0x1D)
GRIS_LABEL = RGBColor(0xA0, 0x90, 0x80)
FILET = RGBColor(0xE5, 0xE0, 0xD8)
FILET_FORT = RGBColor(0xCF, 0xC8, 0xBD)
FONT_CORPS = "Calibri"
FONT_TITRE = "Garamond"

TMPL_ALLOC_CATS = ["Immobilier", "Financier", "Retraite", "Liquidités"]
TMPL_LIQ_CATS = ["Liquide", "Illiquide"]

BAREME_2025 = [
    ("0 %", 0.08, FILET),
    ("11 %", 0.14, FILET_FORT),
    ("30 %", 0.28, GRIS_LABEL),
    ("41 %", 0.26, BRONZE_FORT_RGB),
    ("45 %", 0.24, BORDEAUX),
]


# ---------------------------------------------------------------------------
# 1) Remplissage texte par regex (placeholders [entre crochets])
# ---------------------------------------------------------------------------

def _replace_run_text(run_xml, new_text, target_color=COLOR_BRONZE_FORT):
    esc = new_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    run_xml = re.sub(r"(?<=<a:t>).*?(?=</a:t>)", lambda _m: esc, run_xml, flags=re.S)
    run_xml = re.sub(r'\si="1"', "", run_xml)  # R17 : romain une fois rempli

    def recolor(m):
        if m.group(1) in PLACEHOLDER_COLORS_TO_MIGRATE:
            return f'<a:srgbClr val="{target_color}"/>'
        return m.group(0)

    return re.sub(r'<a:srgbClr val="([0-9A-Fa-f]{6})"/>', recolor, run_xml)


def inject_slide(xml, replacements, mois_annee, target_color=COLOR_BRONZE_FORT):
    xml = xml.replace("[MOIS ANNÉE]", mois_annee)
    runs = list(re.finditer(r"<a:r>(?:(?!</a:r>).)*</a:r>", xml, re.S))
    out, last_end, ri = [], 0, 0
    for m in runs:
        run_xml = m.group(0)
        if "[" in run_xml:
            if ri >= len(replacements):
                raise ValueError(f"plan incomplet : pas assez de valeurs (attendu > {ri})")
            out.append(xml[last_end:m.start()])
            out.append(_replace_run_text(run_xml, replacements[ri], target_color))
            ri += 1
            last_end = m.end()
    out.append(xml[last_end:])
    if ri != len(replacements):
        raise ValueError(f"plan trop long : {ri} runs remplis, {len(replacements)} valeurs fournies")
    return "".join(out)


def inject_slide6_labels(xml, allocation_pct, liquidite_pct):
    """Corrige le bug confirmé le 30/08/2026 : check_charts.py met à jour le
    graphique natif de la slide 6 mais jamais les 6 étiquettes texte '[X %]'
    superposées (objet séparé). Même garde-fou que check_charts.py --fix :
    aucune correction si les catégories ne correspondent pas exactement."""
    if not allocation_pct or not liquidite_pct:
        return xml, "aucune donnée allocation_pct/liquidite_pct — labels non touchés"
    if set(allocation_pct) != set(TMPL_ALLOC_CATS) or set(liquidite_pct) != set(TMPL_LIQ_CATS):
        return xml, "⚠ catégories fournies ≠ catégories du template — labels NON corrigés"

    values = [allocation_pct[c] for c in TMPL_ALLOC_CATS] + [liquidite_pct[c] for c in TMPL_LIQ_CATS]
    runs = list(re.finditer(r"<a:r>(?:(?!</a:r>).)*</a:r>", xml, re.S))
    ph_indices = [i for i, m in enumerate(runs) if "[X %]" in m.group(0)]
    if len(ph_indices) != 6:
        return xml, f"⚠ structure inattendue ({len(ph_indices)} labels au lieu de 6) — abandon"

    out, last_end = [], 0
    for pos, run_i in enumerate(ph_indices):
        m = runs[run_i]
        val = f"{values[pos]:.2f}".replace(".", ",") + " %"
        out.append(xml[last_end:m.start()])
        out.append(_replace_run_text(m.group(0), val))
        last_end = m.end()
    out.append(xml[last_end:])
    return "".join(out), "✓ 6 labels corrigés"


# ---------------------------------------------------------------------------
# 2) Reconstruction native slide 5 (maquette 2 colonnes / famille) + slide 7
#    (barème IR) — via python-pptx, appliqué sur le fichier déjà dézippé par
#    l'étape 1 (donc après le remplissage regex des autres slides, avant
#    ré-empaquetage : les deux passes travaillent sur le même répertoire
#    temporaire pour rester dans un seul appel).
# ---------------------------------------------------------------------------

def _textbox(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             font=FONT_CORPS, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def _rect(slide, l, t, w, h, fill_color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line:
        shp.line.color.rgb = NAVY
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def rebuild_slide5_patrimoine(prs, familles, total_net):
    s5 = list(prs.slides)[4]
    to_remove, frame_bar = [], None
    for sh in s5.shapes:
        if sh.shape_type == 19:  # TABLE
            to_remove.append(sh)
        elif sh.has_text_frame and sh.text_frame.text.strip() == "RÉPARTITION VISUELLE":
            to_remove.append(sh)
        elif (sh.width is not None and sh.width < 20000 and sh.height is not None
              and sh.height > 4000000 and sh.left is not None and sh.left > 7000000):
            # Filet bordeaux vertical décoratif de droite (§5.4) — retiré
            # UNIQUEMENT sur cette slide (demande explicite Marc 30/08),
            # jamais reproduit sur une autre slide du Canvas.
            frame_bar = sh
    for sh in s5.shapes:
        if sh.left is not None and sh.left >= 5900000 and sh.top is not None and 1200000 < sh.top < 4900000:
            if sh not in to_remove:
                to_remove.append(sh)
    if frame_bar is not None:
        to_remove.append(frame_bar)
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)

    def money(x):
        return f"{x:,.0f} €".replace(",", " ")

    margin_left, content_right, gap = 502920, 8460000, 300000
    usable = content_right - margin_left
    col_w = (usable - gap) // 2
    col_xs = [margin_left, margin_left + col_w + gap]

    half = (len(familles) + 1) // 2
    cols = [familles[:half], familles[half:]]

    # CORRECTIF 01/09/2026 (retour Marc, proposition A retenue) : le bloc
    # démarrait toujours à y=1300000 fixe, quel que soit le nombre réel de
    # familles/actifs — avec peu de lignes, un grand vide séparait le bloc de
    # la bannière TOTAL en bas. On centre désormais verticalement le bloc
    # (les 2 colonnes démarrent à la même hauteur, calculée) dans l'espace
    # disponible entre le titre et la bannière — mise en page des colonnes
    # elle-même inchangée.
    FAMILY_HEADER_H, ACTIF_ROW_H, FAMILY_GAP = 210000, 185000, 90000
    CONTENT_TOP, CONTENT_BOTTOM = 1300000, 4460000  # 4460000 = 100000 EMU au-dessus de la bannière (4560000)

    def _col_height(fam_list):
        return sum(FAMILY_HEADER_H + len(fam["actifs"]) * ACTIF_ROW_H + FAMILY_GAP for fam in fam_list)

    block_height = max((_col_height(c) for c in cols), default=0)
    envelope = CONTENT_BOTTOM - CONTENT_TOP
    y = CONTENT_TOP + max(0, envelope - block_height) // 2

    for col_i, fam_list in enumerate(cols):
        cy, cx = y, col_xs[col_i]
        for fam in fam_list:
            sous_total = sum(a["montant"] for a in fam["actifs"])
            _textbox(s5, cx, cy, col_w, 180000, f"{fam['nom']} — {money(sous_total)}", 9,
                     BORDEAUX, bold=True, font=FONT_TITRE)
            cy += 210000
            for a in fam["actifs"]:
                _rect(s5, cx, cy + 155000, col_w, 1000, FILET)
                _textbox(s5, cx, cy, col_w * 0.56, 170000, a["nom"], 8, CORPS)
                prop = a["proprietaire"]
                prop_color = GRIS_LABEL if prop.lower().startswith("commun") else (
                    BORDEAUX if a.get("initiale_bordeaux") else NAVY)
                prop_font = FONT_TITRE if prop_color in (BORDEAUX, NAVY) else FONT_CORPS
                _textbox(s5, cx + col_w * 0.56, cy, col_w * 0.20, 170000, prop, 7.5,
                         prop_color, align=PP_ALIGN.CENTER, font=prop_font)
                _textbox(s5, cx + col_w * 0.76, cy, col_w * 0.24, 170000, money(a["montant"]), 8,
                         CORPS, align=PP_ALIGN.RIGHT)
                cy += 185000
            cy += 90000

    banner_y, banner_w = 4560000, content_right - margin_left
    banner = _rect(s5, margin_left, banner_y, banner_w, 300000, NAVY)
    tf = banner.text_frame
    tf.margin_left, tf.margin_right = Emu(120000), Emu(120000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r1 = tf.paragraphs[0].add_run()
    r1.text = "TOTAL PATRIMOINE NET"
    r1.font.size, r1.font.color.rgb, r1.font.name = Pt(10), CREME, FONT_CORPS
    _textbox(s5, margin_left + banner_w - 2354500, banner_y, 2354500, 300000, money(total_net),
             12, CREME, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_bareme_bar(prs, tmi_label):
    s7 = list(prs.slides)[6]
    x0, y0, w = 502920, 3130000, 3200400
    _textbox(s7, x0, y0, w, 170000, "BARÈME PAR TRANCHE (IR, 1 PART)", 6.5, GRIS_LABEL)
    bar_y, bar_h = y0 + 200000, 190000
    cx = x0
    seg_positions = []
    for label, frac, color in BAREME_2025:
        seg_w = int(w * frac)
        seg_positions.append(cx)
        _rect(s7, cx, bar_y, seg_w, bar_h, color, line=(label == tmi_label))
        cx += seg_w
    label_y = bar_y + bar_h + 30000
    for i, (label, frac, color) in enumerate(BAREME_2025):
        lw = int(w * frac) + 200000
        col = BORDEAUX if label == tmi_label else CORPS
        _textbox(s7, seg_positions[i] - 50000, label_y, lw, 140000, label, 6, col,
                 bold=(label == tmi_label))
    note_y = label_y + 170000
    _textbox(s7, x0, note_y, w, 260000,
             "Seuils 2025 (1 part) — jusqu'à 11 497 € / 29 315 € / 83 823 € / 180 294 € / au-delà",
             5.5, GRIS_LABEL, italic=True)


# ---------------------------------------------------------------------------
# 3) Orchestration
# ---------------------------------------------------------------------------

def _slide_order_by_rid(pptx_path):
    """Retourne la liste des noms de fichier slideN.xml dans l'ordre physique
    réel (celui de <p:sldIdLst> dans presentation.xml), en résolvant chaque
    r:id via presentation.xml.rels. Fonctionne sur un .pptx classique (pas
    un dossier extrait) — utilisé avant ET après le passage python-pptx pour
    comparer les deux états."""
    with zipfile.ZipFile(pptx_path) as z:
        pres_xml = z.read("ppt/presentation.xml").decode("utf-8")
        rels_xml = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")

    rid_to_target = dict(re.findall(
        r'Id="(rId\d+)"[^>]*Target="slides/([^"]+)"', rels_xml))
    ordered_rids = re.findall(r'<p:sldId[^>]*r:id="(rId\d+)"', pres_xml)
    return [rid_to_target[r] for r in ordered_rids if r in rid_to_target]


def restore_slide_filenames(pptx_path, expected_order):
    """CORRECTIF 31/08/2026 (crash-test consultant, dossier réel) — bug
    critique : toute sauvegarde python-pptx (`Presentation(...).save(...)`)
    renumérote les fichiers ppt/slides/slideN.xml selon leur POSITION
    physique (1..N sans trous), quel que soit leur nom avant l'ouverture.
    `template_positions.json`, `nb_elements_reels` et le routage couleur de
    ce script identifient pourtant les slides par leur nom de fichier
    D'ORIGINE (12=Projets, 13=Moyens, 15=Synthèse...) — jamais par position.
    Conséquence observée en test réel : après la reconstruction native de la
    slide 5 (patrimoine) et/ou de la barre de barème de la slide 7 (qui
    imposent un aller-retour python-pptx), `remove_unused_slots.py` ciblait
    la mauvaise slide et ne supprimait rien (slot factice resté visible dans
    le deck livré au client).

    Cette fonction compare l'ordre `expected_order` (noms de fichier corrects,
    capturés AVANT l'ouverture python-pptx) à l'ordre réel après sauvegarde,
    et renomme chaque partie (slideN.xml, son _rels, les cibles dans
    presentation.xml.rels et [Content_Types].xml, et les .rels des
    notesSlides qui référencent la slide) pour restaurer les noms d'origine.
    Idempotent : si aucune renumérotation n'a eu lieu, ne modifie rien."""
    tmp = pptx_path + ".__tmp_restore"
    if os.path.exists(tmp):
        shutil.rmtree(tmp)
    os.makedirs(tmp)
    with zipfile.ZipFile(pptx_path) as z:
        z.extractall(tmp)

    actual_order = _slide_order_by_rid(pptx_path)
    if actual_order == expected_order:
        shutil.rmtree(tmp)
        return False  # rien à corriger

    if len(actual_order) != len(expected_order):
        shutil.rmtree(tmp)
        raise ValueError(
            f"restore_slide_filenames : nombre de slides incohérent "
            f"({len(actual_order)} vs {len(expected_order)} attendues) — "
            f"correction automatique impossible, à traiter manuellement."
        )

    # rename_map : nom ACTUEL (après save) -> nom ATTENDU (avant save),
    # à la même position physique. On passe par un préfixe temporaire pour
    # ne jamais écraser un fichier cible avant de l'avoir déplacé.
    rename_map = dict(zip(actual_order, expected_order))
    slides_dir = os.path.join(tmp, "ppt", "slides")
    rels_dir = os.path.join(slides_dir, "_rels")

    for actual in rename_map:
        stem = actual.replace(".xml", "")
        for p in (os.path.join(slides_dir, actual),
                  os.path.join(rels_dir, f"{stem}.xml.rels")):
            if os.path.exists(p):
                os.rename(p, p + ".__tmp")
    for actual, expected in rename_map.items():
        stem = actual.replace(".xml", "")
        new_stem = expected.replace(".xml", "")
        p_xml = os.path.join(slides_dir, actual) + ".__tmp"
        if os.path.exists(p_xml):
            os.rename(p_xml, os.path.join(slides_dir, expected))
        p_rels = os.path.join(rels_dir, f"{stem}.xml.rels") + ".__tmp"
        if os.path.exists(p_rels):
            os.rename(p_rels, os.path.join(rels_dir, f"{new_stem}.xml.rels"))

    def _rewrite(path, mapping):
        """Une seule passe via regex + callback : chaque correspondance est
        résolue indépendamment contre le texte ORIGINAL, jamais contre le
        résultat d'une substitution précédente — élimine tout risque de
        cascade (ex. 10->11 puis 11->12 qui rattraperait le 10 fraîchement
        converti). Bug réel trouvé le 31/08/2026 en testant cette fonction
        elle-même sur une 2e édition du Canvas Master : des substitutions
        successives naïves (str.replace en boucle) avaient fait converger
        plusieurs rId vers le même fichier cible (slide19.xml en x6)."""
        data = open(path, encoding="utf-8").read()
        pattern = re.compile("|".join(re.escape(k) for k in mapping))
        data = pattern.sub(lambda m: mapping[m.group(0)], data)
        open(path, "w", encoding="utf-8").write(data)

    slide_map = {f'Target="slides/{a}"': f'Target="slides/{e}"'
                 for a, e in rename_map.items()}
    _rewrite(os.path.join(tmp, "ppt", "_rels", "presentation.xml.rels"), slide_map)
    ct_map = {f'/ppt/slides/{a}"': f'/ppt/slides/{e}"'
              for a, e in rename_map.items()}
    _rewrite(os.path.join(tmp, "[Content_Types].xml"), ct_map)

    notes_rels_dir = os.path.join(tmp, "ppt", "notesSlides", "_rels")
    if os.path.isdir(notes_rels_dir):
        notes_map = {f'Target="../slides/{a}"': f'Target="../slides/{e}"'
                     for a, e in rename_map.items()}
        for fn in os.listdir(notes_rels_dir):
            _rewrite(os.path.join(notes_rels_dir, fn), notes_map)

    fixed = pptx_path + ".__tmp_fixed.pptx"
    if os.path.exists(fixed):
        os.remove(fixed)
    with zipfile.ZipFile(fixed, "w", zipfile.ZIP_DEFLATED) as zf:
        for root, _, files in os.walk(tmp):
            for f in files:
                full = os.path.join(root, f)
                zf.write(full, os.path.relpath(full, tmp))
    shutil.rmtree(tmp)
    os.replace(fixed, pptx_path)
    print(f"  ⚠ renumérotation python-pptx détectée et corrigée automatiquement "
          f"({len(rename_map)} fichier(s) renommé(s))")
    return True


def show_fields(canvas_path, slide_num):
    """CORRECTIF 31/08/2026 (crash-test consultant, dossier réel F.D./I.D.) —
    3 champs de la slide 7 (RATIO D'ENDETTEMENT, MARGE AVANT TRANCHE SUP.,
    REVENUS NETS MENSUELS) ont reçu des valeurs qui ne correspondaient pas à
    leur libellé réel : le format `replacements` est une liste positionnelle
    pure (l'ordre d'apparition dans le XML, invisible sans l'ouvrir), et rien
    n'empêchait de deviner cet ordre au lieu de le vérifier.

    Cette commande affiche, pour une slide donnée, la séquence exacte
    attendue par `inject_slide()` — index, libellé le plus proche trouvé
    juste avant chaque run à crochets, et texte du placeholder — à consulter
    AVANT de préparer `replacements[slide_num]`, plutôt que de deviner
    l'ordre. Ne modifie rien, ne nécessite pas de plan.json.
    """
    with zipfile.ZipFile(canvas_path) as z:
        xml = z.read(f"ppt/slides/slide{slide_num}.xml").decode("utf-8")
    xml_for_scan = xml.replace("[MOIS ANNÉE]", "X")

    all_texts = list(re.finditer(r"<a:t>([^<]*)</a:t>", xml_for_scan))
    bracket_positions = {m.start() for m in all_texts if "[" in m.group(1)}

    print(f"\nSlide {slide_num} — {len(bracket_positions)} champ(s) attendu(s) dans "
          f"cet ordre exact pour replacements[\"{slide_num}\"] :\n")
    idx = 0
    last_label = None
    for m in all_texts:
        text = m.group(1)
        if m.start() in bracket_positions:
            idx += 1
            label = last_label if last_label else "(aucun libellé détecté avant ce champ)"
            print(f"  [{idx}] libellé le plus proche : {label!r:55}  placeholder : {text!r}")
        elif text.strip():
            last_label = text.strip()
    print()


def personalize(canvas_in, plan, out_path):
    tmp_dir = out_path + ".__tmp_personalize"
    if os.path.exists(tmp_dir):
        shutil.rmtree(tmp_dir)
    os.makedirs(tmp_dir)
    with zipfile.ZipFile(canvas_in) as z:
        z.extractall(tmp_dir)

    mois_annee = plan.get("mois_annee", "")
    replacements = dict(plan.get("replacements", {}))
    if "5" in replacements or 5 in replacements:
        raise ValueError("plan invalide : la clé '5' ne doit jamais être fournie dans "
                          "'replacements' — la slide 5 est reconstruite nativement (patrimoine).")

    def _migrate_remaining_bronze_clair(xml, target_color):
        """CORRECTIF 01/09/2026 (crash-test 3 dossiers fictifs) : le remplissage par
        placeholder ne recolore QUE les runs effectivement substitués. Deux cas
        échappaient donc à R16 :
          1. le token '[MOIS ANNÉE]' (remplacement de chaîne brut, en amont de la
             boucle par run) — gap déjà documenté ;
          2. tout texte FIXE du Canvas (ex. le sous-titre 'Stratégie patrimoniale'
             en couverture) qui n'est jamais un placeholder '[...]' mais reste
             néanmoins stylé en bronze clair 8E5B3F dans le fichier maître — jamais
             documenté avant ce crash-test.
        Passe finale, globale à la slide : tout 8E5B3F restant migre vers la
        couleur cible de CETTE slide (bronze fort ou taupe selon color_for_slide),
        conformément à l'esprit de R16 ('toute occurrence... migre au remplissage')."""
        return re.sub(
            r'<a:srgbClr val="8E5B3F"/>',
            f'<a:srgbClr val="{target_color}"/>',
            xml,
        )

    for slide_num, values in replacements.items():
        path = os.path.join(tmp_dir, "ppt", "slides", f"slide{slide_num}.xml")
        xml = open(path, encoding="utf-8").read()
        target_color = color_for_slide(slide_num)
        xml = inject_slide(xml, values, mois_annee, target_color)
        xml = _migrate_remaining_bronze_clair(xml, target_color)
        open(path, "w", encoding="utf-8").write(xml)

    path6 = os.path.join(tmp_dir, "ppt", "slides", "slide6.xml")
    xml6 = open(path6, encoding="utf-8").read().replace("[MOIS ANNÉE]", mois_annee)
    xml6, status6 = inject_slide6_labels(xml6, plan.get("allocation_pct"), plan.get("liquidite_pct"))
    xml6 = _migrate_remaining_bronze_clair(xml6, color_for_slide("6"))
    open(path6, "w", encoding="utf-8").write(xml6)
    print(f"  slide6 (répartition) : {status6}")

    repacked = out_path + ".__tmp_repacked.pptx"
    if os.path.exists(repacked):
        os.remove(repacked)
    with zipfile.ZipFile(repacked, "w", zipfile.ZIP_DEFLATED) as zf:
        for root, _, files in os.walk(tmp_dir):
            for f in files:
                full = os.path.join(root, f)
                zf.write(full, os.path.relpath(full, tmp_dir))
    shutil.rmtree(tmp_dir)

    # Capturé AVANT toute ouverture python-pptx : c'est le SEUL moment où les
    # noms de fichier sont garantis corrects (cf. restore_slide_filenames).
    expected_order = _slide_order_by_rid(repacked)

    prs = Presentation(repacked)
    patrimoine = plan.get("patrimoine")
    if patrimoine:
        rebuild_slide5_patrimoine(prs, patrimoine["familles"], patrimoine["total_net"])
    else:
        print("  ⚠ pas de bloc 'patrimoine' dans le plan — slide 5 laissée au Canvas d'origine "
              "(non recommandé, cf. §5.2bis)")
    tmi_label = plan.get("tmi_label")
    if tmi_label:
        add_bareme_bar(prs, tmi_label)
    prs.save(out_path)
    os.remove(repacked)
    restore_slide_filenames(out_path, expected_order)
    print(f"OK -> {out_path}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--canvas", required=True)
    ap.add_argument("--plan", required=False, help="fichier JSON, cf. docstring "
                    "(non requis avec --show-fields)")
    ap.add_argument("--out", required=False)
    ap.add_argument("--show-fields", metavar="SLIDE_NUM", default=None,
                     help="n'écrit rien : affiche l'ordre exact des champs attendus "
                          "pour cette slide (à consulter avant de préparer 'replacements', "
                          "cf. bug de mapping constaté le 31/08/2026)")
    args = ap.parse_args()
    if args.show_fields:
        show_fields(args.canvas, args.show_fields)
        raise SystemExit(0)
    if not args.plan or not args.out:
        raise SystemExit("--plan et --out sont requis (sauf en mode --show-fields)")
    plan_data = json.load(open(args.plan, encoding="utf-8"))
    personalize(args.canvas, plan_data, args.out)
