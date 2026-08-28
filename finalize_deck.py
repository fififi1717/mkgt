#!/usr/bin/env python3
"""
Affilior — Finalisation post-remplissage d'un dossier (Étape 5, après le
remplissage de contenu qui suit assemble.py).

Chaîne 4 contrôles/corrections (mise à jour v4.14, ajout de check_charts issu
d'une session parallèle du 28/08/2026 — auparavant 3 contrôles seulement,
issus du crash-test du 27/08/2026), dans cet ordre (chacun opère sur le
fichier produit par le précédent) :

  1. remove_unused_slots  — supprime réellement les formes des slots vides
                             (règle 5.2quater) plutôt que de vider leur texte.
  2. shrink_oversized_text — réduit automatiquement (paliers ~30%, 2 max) le
                             texte qui déborde des zones "gros chiffre" ;
                             signale ⚠ ce qui reste illisible plutôt que de
                             deviner davantage.
  3. check_charts          — contrôle NON BLOQUANT (+ correction optionnelle
                             via --fix-charts) des graphiques natifs (donuts) :
                             assemble.py ne touche jamais aux séries de
                             graphique, seuls les placeholders texte/tableaux
                             sont remplis — un deck peut donc être livré avec
                             les donuts encore aux valeurs de démo du template
                             pendant que le texte autour affiche les vrais
                             chiffres (bug confirmé crash-test 28/08/2026).
  4. check_montants        — contrôle NON BLOQUANT de traçabilité : tout
                             montant du deck introuvable dans les sources
                             déclarées (chiffres_source) est signalé, jamais
                             bloqué (la validation métier reste à l'associé
                             IP — cf. décision du 27/08/2026).

Rien ici ne remplace la QA visuelle (planche contact, §5.6) ni la validation
associé IP (R4). Ce sont des filets de sécurité mécaniques pour les erreurs
que les runs précédents ont montré comme non détectées par validate.py seul.
Voir aussi check_template_residue.py (script séparé, non chaîné ici) pour
2 défauts supplémentaires du même ordre (couleur placeholder non corrigée,
notes de conception internes visibles) — à exécuter en plus, pas à la place.

Usage :
    python3 finalize_deck.py deck.pptx --out deck_final.pptx \
        --template-positions template_positions.json \
        --dossier dossier_client.json [--fix-charts]

template_positions.json est GÉNÉRIQUE — commun à tous les clients, ne change
que si Cans_Mstr.pptx change de version. Contient les coordonnées XML des
slots variables et la liste fixe des slides personnalisées.

dossier_client.json est PROPRE À CE CLIENT — à reconstruire à chaque dossier :
{
  "nb_elements_reels": {"12": 2, "13": 2, "15": 2},
  "chiffres_source": [5880000, 225783, ...],
  "allocation_pct": {"Immobilier": 45, "Financier": 38, "Retraite": 12, "Liquidités": 5},
  "liquidite_pct": {"Liquide": 20, "Illiquide": 80}
}
allocation_pct et liquidite_pct sont OPTIONNELS (ajoutés v4.14, cf. check_charts.py) :
sans eux, l'étape 3/4 se contente de signaler ⚠ si un graphique est resté en
valeurs de démo, sans pouvoir le corriger automatiquement même avec --fix-charts.
"""
import argparse
import json
import shutil
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
import remove_unused_slots
import shrink_oversized_text
import check_montants
import check_charts
# select_slots.py fusionné dans remove_unused_slots.py le 28/08/2026 (nettoyage dépôt) —
# compute_slots_a_supprimer y est désormais définie directement.


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("src")
    ap.add_argument("--out", required=True)
    ap.add_argument("--template-positions", required=True,
                     help="Référentiel générique (identique pour tous les clients, "
                          "cf. template_positions.json)")
    ap.add_argument("--dossier", required=True,
                     help="Fichier propre à CE client (nb_elements_reels + chiffres_source)")
    ap.add_argument("--tolerance", type=int, default=0)
    ap.add_argument("--fix-charts", action="store_true",
                     help="Corrige automatiquement les graphiques natifs (donuts) restés aux "
                          "valeurs de démo, si allocation_pct/liquidite_pct sont fournis dans "
                          "--dossier et sans ambiguïté (cf. check_charts.py).")
    args = ap.parse_args()

    template_positions = json.load(open(args.template_positions, encoding="utf-8"))
    dossier = json.load(open(args.dossier, encoding="utf-8"))

    slots = remove_unused_slots.compute_slots_a_supprimer(
        template_positions, dossier.get("nb_elements_reels", {})
    )
    slides_check = template_positions["slides_personnalisees"]
    chiffres_source = set(dossier.get("chiffres_source", []))

    with tempfile.TemporaryDirectory() as tmp:
        step1 = str(Path(tmp) / "step1_slots.pptx")
        step2 = str(Path(tmp) / "step2_shrink.pptx")

        print("=== [1/4] Suppression des slots vides (5.2quater) ===")
        if slots:
            remove_unused_slots.apply(args.src, step1, slots)
        else:
            shutil.copy(args.src, step1)
            print("Aucun slot déclaré à supprimer — étape ignorée.")

        print("\n=== [2/4] Réduction des zones en débordement ===")
        report_shrink = shrink_oversized_text.apply(step1, step2, slides_check)

        step3 = str(Path(tmp) / "step3_charts.pptx")
        print("\n=== [3/4] Contrôle des graphiques natifs (donuts, non bloquant) ===")
        chart_alerts, chart_fixes = check_charts.check(step2, dossier)
        chart_fixed_slides = set()
        if chart_alerts:
            for slide_n, name, cats, vals, msg in chart_alerts:
                print(f"  slide {slide_n:<3} série « {name} » {dict(zip(cats, vals))}")
                print(f"    {msg}")
        else:
            print("✓ Aucun graphique laissé aux valeurs de démonstration du template.")

        if args.fix_charts and chart_fixes:
            import check_charts as _cc_mod
            from pptx import Presentation as _Presentation
            from pptx.chart.data import CategoryChartData as _CCD
            prs_charts = _Presentation(step2)
            expected = {"Allocation": dossier.get("allocation_pct"),
                        "Liquidité": dossier.get("liquidite_pct")}
            for slide_n, slide in enumerate(prs_charts.slides, start=1):
                for shape in slide.shapes:
                    if not shape.has_chart:
                        continue
                    chart = shape.chart
                    plot = chart.plots[0]
                    categories = list(plot.categories)
                    for series in plot.series:
                        values = list(series.values)
                        exp = expected.get(series.name)
                        if exp and _cc_mod._matches_demo(categories, values) \
                                and set(exp.keys()) == set(categories) \
                                and abs(sum(exp.values()) - 100) <= 0.5:
                            cd = _CCD()
                            cd.categories = categories
                            cd.add_series(series.name, [exp[c] for c in categories])
                            chart.replace_data(cd)
                            chart_fixed_slides.add(slide_n)
            prs_charts.save(step3)
            print(f"✓ graphique(s) corrigé(s) automatiquement sur la/les slide(s) {sorted(chart_fixed_slides)}.")
        else:
            shutil.copy(step2, step3)
            if chart_fixes and not args.fix_charts:
                print("  (correction disponible mais --fix-charts non passé — donuts laissés tels quels)")

        shutil.copy(step3, args.out)
        print(f"\n-> Fichier finalisé : {args.out}")

        print("\n=== [4/4] Contrôle de traçabilité des montants (non bloquant) ===")
        allowed_slides = {f"ppt/slides/slide{n}.xml" for n in slides_check}
        found = [
            (s, r, v)
            for (s, r, v) in check_montants.extract_amounts_from_pptx(args.out)
            if s in allowed_slides
        ]
        seen = set()
        orphans = []
        for slide, raw, val in found:
            key = (slide, val)
            if key in seen:
                continue
            seen.add(key)
            if not any(abs(val - r) <= args.tolerance for r in chiffres_source):
                orphans.append((slide, raw, val))

        print(f"Montants distincts contrôlés : {len(seen)} | Sources tracées : {len(chiffres_source)}")
        if orphans:
            print(f"⚠ {len(orphans)} montant(s) non tracé(s) — à vérifier avant diffusion :")
            for slide, raw, val in orphans:
                print(f"  ⚠ {slide.replace('ppt/slides/', ''):15s} {raw:>15s}")
        else:
            print("✓ Tous les montants contrôlés sont tracés à une source déclarée.")

        print("\n=== Résumé pour le message de livraison (§5.7) ===")
        unresolved = [r for r in report_shrink if not r[4]]
        if unresolved:
            print(f"⚠ {len(unresolved)} zone(s) toujours en débordement après réduction automatique "
                  f"— vérification visuelle obligatoire avant envoi.")
        if orphans:
            print(f"⚠ {len(orphans)} montant(s) non tracé(s) à une source — à confirmer.")
        if chart_alerts and not chart_fixed_slides:
            print(f"⚠ {len(chart_alerts)} graphique(s) natif(s) resté(s) en valeurs de démo — non corrigé(s).")
        if not unresolved and not orphans and not (chart_alerts and not chart_fixed_slides):
            print("✓ Aucun point mécanique résiduel. QA visuelle (§5.6) et validation associé IP (R4) restent requises.")

        # --- Périmètre restreint pour la QA visuelle (§5.6) ---
        # Seules les slides effectivement touchées par ce run ont besoin d'être
        # rendues en image pour vérification humaine — pas les 28 slides du
        # deck. Évite un rendu soffice/pdftoppm complet à chaque run.
        import re as _re
        touched_slides = set(slides_check) | chart_fixed_slides | {
            int(_re.search(r"slide(\d+)\.xml$", s).group(1)) for s, _, _ in found
        }
        if touched_slides:
            page_list = ",".join(str(n) for n in sorted(touched_slides))
            print(f"\n=== Périmètre QA visuelle recommandé (§5.6) ===")
            print(f"Slides à rendre en image (pas le deck entier) : {page_list}")
            print(f"  soffice --headless --convert-to pdf {args.out}")
            print(f"  pdftoppm -jpeg -r 150 -f <min> -l <max> {Path(args.out).stem}.pdf apercu")
            print("  (ou extraire page par page avec -f N -l N pour chaque slide de la liste ci-dessus)")


if __name__ == "__main__":
    main()
